"""Utility to regenerate PPT decks in multiple AI-inspired styles."""
from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE_DIR = Path(__file__).resolve().parent
DATA_PATH = BASE_DIR / "data" / "presentation_outline.json"
OUTPUT_DIR = BASE_DIR / "output"

StyleDict = Dict[str, object]

STYLE_PRESETS: Dict[str, StyleDict] = {
    "original": {
        "title": "Original Baseline",
        "filename": "smart_campus_original.pptx",
        "palette": {
            "background": (255, 255, 255),
            "title": (20, 20, 20),
            "body": (80, 80, 80),
            "accent": (23, 96, 167),
            "accent_alt": (231, 242, 250),
            "card_border": (200, 200, 200),
        },
        "font": {"title": "Calibri", "body": "Calibri"},
        "decor": "none",
    },
    "pulse_neon": {
        "title": "Pulse Neon (AI Redesign)",
        "filename": "smart_campus_pulse_neon.pptx",
        "palette": {
            "background": (5, 8, 28),
            "title": (250, 250, 250),
            "body": (199, 205, 255),
            "accent": (120, 82, 248),
            "accent_alt": (19, 196, 223),
            "card_border": (80, 90, 170),
        },
        "font": {"title": "Montserrat", "body": "Inter"},
        "decor": "blobs",
    },
    "zen_canvas": {
        "title": "Zen Canvas (AI Redesign)",
        "filename": "smart_campus_zen_canvas.pptx",
        "palette": {
            "background": (248, 247, 242),
            "title": (24, 36, 42),
            "body": (66, 87, 99),
            "accent": (134, 153, 110),
            "accent_alt": (200, 211, 167),
            "card_border": (210, 206, 192),
        },
        "font": {"title": "Source Sans Pro", "body": "Source Sans Pro"},
        "decor": "strokes",
    },
}


def load_outline() -> Dict[str, object]:
    with DATA_PATH.open("r", encoding="utf-8") as fh:
        return json.load(fh)


def apply_background(slide, style: StyleDict) -> None:
    fill = slide.background.fill
    fill.solid()
    r, g, b = style["palette"]["background"]
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_title_box(slide, text: str, style: StyleDict, top: Inches = Inches(0.7)) -> None:
    tx_box = slide.shapes.add_textbox(Inches(0.7), top, Inches(11.5), Inches(1.8))
    tf = tx_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(46)
    p.font.name = style["font"]["title"]
    r, g, b = style["palette"]["title"]
    p.font.color.rgb = RGBColor(r, g, b)


def add_subtitle(slide, text: str, style: StyleDict, top_offset: float = 2.0) -> None:
    tx_box = slide.shapes.add_textbox(Inches(0.7), Inches(top_offset), Inches(9.5), Inches(1.2))
    tf = tx_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.font.name = style["font"]["body"]
    r, g, b = style["palette"]["body"]
    p.font.color.rgb = RGBColor(r, g, b)


def add_badge(slide, text: str, style: StyleDict) -> None:
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.9), Inches(3.0), Inches(0.6))
    fill = badge.fill
    fill.solid()
    r, g, b = style["palette"]["accent"]
    fill.fore_color.rgb = RGBColor(r, g, b)
    badge.line.color.rgb = RGBColor(r, g, b)
    badge.text_frame.text = text
    p = badge.text_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = style["font"]["body"]
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def add_bullets(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    if desc := data.get("description"):
        add_subtitle(slide, desc, style, top_offset=1.9)
    tx_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.8), Inches(10.8), Inches(3.8))
    tf = tx_box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(data.get("bullets", [])):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.name = style["font"]["body"]
        r, g, b = style["palette"]["body"]
        p.font.color.rgb = RGBColor(r, g, b)
        p.level = 0


def add_split_highlight(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    left_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.3), Inches(5.5), Inches(3.8))
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.6), Inches(2.3), Inches(5.5), Inches(3.8))
    for idx, box in enumerate((left_box, right_box)):
        fill = box.fill
        fill.solid()
        accent_key = "accent" if idx == 1 else "accent_alt"
        r, g, b = style["palette"][accent_key]
        fill.fore_color.rgb = RGBColor(r, g, b)
        box.line.color.rgb = RGBColor(r, g, b)
    for block, payload in zip((left_box, right_box), (data["left"], data["right"])):
        tf = block.text_frame
        tf.clear()
        heading = tf.paragraphs[0]
        heading.text = payload["heading"]
        heading.font.size = Pt(22)
        heading.font.bold = True
        heading.font.color.rgb = RGBColor(255, 255, 255)
        heading.font.name = style["font"]["body"]
        for item in payload["items"]:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.name = style["font"]["body"]
            p.font.color.rgb = RGBColor(255, 255, 255)


def add_architecture(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    add_subtitle(slide, data.get("description", ""), style, top_offset=1.9)
    nodes: List[Dict[str, str]] = data.get("nodes", [])
    width = Inches(11.2)
    step_width = width / max(1, len(nodes))
    top = Inches(3.0)
    for idx, node in enumerate(nodes):
        left = Inches(0.7) + step_width * idx
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, step_width - Inches(0.2), Inches(2.4))
        fill = box.fill
        fill.solid()
        r, g, b = style["palette"]["accent_alt"] if idx % 2 else style["palette"]["accent"]
        fill.fore_color.rgb = RGBColor(r, g, b)
        box.line.color.rgb = RGBColor(r, g, b)
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        title_p = tf.paragraphs[0]
        title_p.text = node["name"]
        title_p.font.size = Pt(20)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.name = style["font"]["body"]
        body = tf.add_paragraph()
        body.text = node["detail"]
        body.font.size = Pt(16)
        body.font.color.rgb = RGBColor(235, 235, 235)
        body.font.name = style["font"]["body"]


def add_metrics(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    cards = data.get("data_points", [])
    card_width = Inches(3.5)
    gap = Inches(0.4)
    top = Inches(2.4)
    for idx, card in enumerate(cards):
        left = Inches(0.7) + idx * (card_width + gap)
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_width, Inches(2.2))
        rect.fill.solid()
        r, g, b = style["palette"]["accent_alt"]
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        r2, g2, b2 = style["palette"]["card_border"]
        rect.line.color.rgb = RGBColor(r2, g2, b2)
        tf = rect.text_frame
        tf.clear()
        value = tf.paragraphs[0]
        value.text = card["value"]
        value.font.size = Pt(36)
        value.font.bold = True
        value.font.name = style["font"]["title"]
        value.font.color.rgb = RGBColor(*style["palette"]["title"])
        label = tf.add_paragraph()
        label.text = card["label"]
        label.font.size = Pt(18)
        label.font.bold = True
        label.font.name = style["font"]["body"]
        label.font.color.rgb = RGBColor(*style["palette"]["body"])
        detail = tf.add_paragraph()
        detail.text = card.get("detail", "")
        detail.font.size = Pt(14)
        detail.font.name = style["font"]["body"]
        detail.font.color.rgb = RGBColor(*style["palette"]["body"])
    if quote := data.get("quote"):
        quote_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(10.8), Inches(1.2))
        tf = quote_box.text_frame
        tf.text = quote
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.name = style["font"]["body"]
        p.font.color.rgb = RGBColor(*style["palette"]["body"])


def add_timeline(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    milestones = data.get("milestones", [])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(3.4), Inches(10.6), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(*style["palette"]["body"])
    line.line.fill.background()
    for idx, milestone in enumerate(milestones):
        left = Inches(1.0) + Inches(3.4) * idx
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, Inches(3.1), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RGBColor(*style["palette"]["accent"])
        circle.line.color.rgb = RGBColor(*style["palette"]["accent"])
        label_box = slide.shapes.add_textbox(left - Inches(0.4), Inches(3.8), Inches(2.0), Inches(1.2))
        tf = label_box.text_frame
        tf.word_wrap = True
        tf.text = milestone["label"]
        lab = tf.paragraphs[0]
        lab.font.bold = True
        lab.font.size = Pt(18)
        lab.font.name = style["font"]["body"]
        lab.font.color.rgb = RGBColor(*style["palette"]["title"])
        detail = tf.add_paragraph()
        detail.text = milestone["detail"]
        detail.font.size = Pt(15)
        detail.font.name = style["font"]["body"]
        detail.font.color.rgb = RGBColor(*style["palette"]["body"])


def add_cta(slide, data: Dict[str, object], style: StyleDict) -> None:
    add_title_box(slide, data["title"], style)
    tx_box = slide.shapes.add_textbox(Inches(0.85), Inches(2.3), Inches(10.8), Inches(2.5))
    tf = tx_box.text_frame
    tf.clear()
    for idx, bullet in enumerate(data.get("bullets", [])):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"→ {bullet}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = style["font"]["body"]
        p.font.color.rgb = RGBColor(*style["palette"]["body"])
    if footer := data.get("footer"):
        footer_box = slide.shapes.add_textbox(Inches(0.85), Inches(4.6), Inches(8.0), Inches(0.8))
        footer_box.text_frame.text = footer
        p = footer_box.text_frame.paragraphs[0]
        p.font.size = Pt(16)
        p.font.name = style["font"]["body"]
        r, g, b = style["palette"]["body"]
        p.font.color.rgb = RGBColor(r, g, b)


def add_decor(slide, style: StyleDict, index: int) -> None:
    decor_type = style.get("decor", "none")
    if decor_type == "blobs":
        blob = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.5), Inches(-0.8), Inches(4.5), Inches(4.5))
        blob.fill.solid()
        r, g, b = style["palette"]["accent"]
        blob.fill.fore_color.rgb = RGBColor(r, g, b)
        blob.fill.fore_color.brightness = 0.4
        blob.line.fill.background()
        blob2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-1.2), Inches(4.2), Inches(5.2), Inches(3.3))
        blob2.fill.solid()
        r2, g2, b2 = style["palette"]["accent_alt"]
        blob2.fill.fore_color.rgb = RGBColor(r2, g2, b2)
        blob2.fill.fore_color.brightness = 0.3
        blob2.line.fill.background()
    elif decor_type == "strokes":
        height = Inches(0.12)
        for offset in range(4):
            bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.7 + 0.3 * offset), Inches(1.6 + offset * 0.5), height)
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*style["palette"]["accent"])
            bar.line.fill.background()


def create_slide(prs: Presentation, slide_data: Dict[str, object], style: StyleDict, index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, style)
    if style.get("decor") != "none":
        add_decor(slide, style, index)
    slide_type = slide_data["type"]
    if slide_type == "title":
        add_title_box(slide, slide_data["title"], style)
        if subtitle := slide_data.get("subtitle"):
            add_subtitle(slide, subtitle, style)
        if badge := slide_data.get("badge"):
            add_badge(slide, badge, style)
    elif slide_type == "bullets":
        add_bullets(slide, slide_data, style)
    elif slide_type == "split_highlight":
        add_split_highlight(slide, slide_data, style)
    elif slide_type == "architecture":
        add_architecture(slide, slide_data, style)
    elif slide_type == "metrics":
        add_metrics(slide, slide_data, style)
    elif slide_type == "timeline":
        add_timeline(slide, slide_data, style)
    elif slide_type == "cta":
        add_cta(slide, slide_data, style)
    else:
        add_bullets(slide, slide_data, style)


def build_presentation(style_key: str, outline: Dict[str, object]) -> Path:
    style = STYLE_PRESETS[style_key]
    prs = Presentation()
    prs.slide_height = Inches(7.5)
    prs.slide_width = Inches(13.33)
    for idx, slide in enumerate(outline["slides"]):
        create_slide(prs, slide, style, idx)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / style["filename"]
    prs.save(out_path)
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate PPT decks in multiple styles")
    parser.add_argument("--style", choices=list(STYLE_PRESETS.keys()) + ["all"], default="all", help="Specific style id or 'all'")
    args = parser.parse_args()
    outline = load_outline()
    style_keys = STYLE_PRESETS.keys() if args.style == "all" else [args.style]
    for key in style_keys:
        path = build_presentation(key, outline)
        print(f"Generated {STYLE_PRESETS[key]['title']} → {path}")


if __name__ == "__main__":
    main()
